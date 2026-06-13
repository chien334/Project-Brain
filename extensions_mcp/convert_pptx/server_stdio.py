from mcp.server.fastmcp import FastMCP
from pydantic import BaseModel, Field
from typing import Optional
from pptx import Presentation
from pptx.util import Inches
import os
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.parts.image import ImagePart

mcp = FastMCP("pptx_mcp")

class PptxFilePath(BaseModel):
    file_path: str = Field(..., description="Path to the PPTX file (e.g., 'C:/data/file.pptx')")

class ExtractImagesInput(BaseModel):
    file_path: str = Field(..., description="Path to the PPTX file.")
    output_dir: str = Field(..., description="Directory to save the extracted images.")

class CreatePptxInput(BaseModel):
    file_path: str = Field(..., description="Path to save the new PPTX file.")
    title: str = Field(..., description="The title of the presentation.")
    content: str = Field(..., description="The content for the first slide.")

@mcp.tool(
    name="pptx_extract_text",
    annotations={
        "readOnlyHint": True,
    }
)
async def extract_text(params: PptxFilePath) -> str:
    """Extracts all text from a PPTX file."""
    try:
        prs = Presentation(params.file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            full_text.append(run.text)
        return '\n'.join(full_text)
    except Exception as e:
        return f"Error extracting text from {params.file_path}: {e}"

@mcp.tool(
    name="pptx_create_presentation",
    annotations={
        "destructiveHint": False, # Not destructive as it creates a new file
    }
)
async def create_presentation(params: CreatePptxInput) -> str:
    """Creates a new PPTX file with a title and content."""
    try:
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = params.title
        subtitle.text = params.content

        prs.save(params.file_path)
        return f"Successfully created PPTX file at {params.file_path}"
    except Exception as e:
        return f"Error creating PPTX file at {params.file_path}: {e}"

@mcp.tool(
    name="pptx_extract_images",
    annotations={
        "readOnlyHint": True,
    }
)
async def extract_images(params: ExtractImagesInput) -> str:
    """Extracts all images from a PPTX file and saves them to a directory."""
    try:
        if not os.path.exists(params.output_dir):
            os.makedirs(params.output_dir)

        prs = Presentation(params.file_path)
        image_paths = []
        
        image_parts = [part for part in prs.iter_parts() if isinstance(part, ImagePart)]

        for i, image_part in enumerate(image_parts):
            image_bytes = image_part.blob
            image_ext = image_part.ext
            image_filename = f"image_{i}.{image_ext}"
            image_path = os.path.join(params.output_dir, image_filename)
            with open(image_path, "wb") as f:
                f.write(image_bytes)
            image_paths.append(image_path)
            
        return f"Successfully extracted {len(image_paths)} images to {params.output_dir}: {image_paths}"
    except Exception as e:
        return f"Error extracting images from {params.file_path}: {e}"

@mcp.tool(
    name="pptx_to_markdown",
    annotations={
        "readOnlyHint": True,
    }
)
async def to_markdown(params: PptxFilePath) -> str:
    """Converts a PPTX file to a Markdown document."""
    try:
        prs = Presentation(params.file_path)
        markdown_text = ""
        for i, slide in enumerate(prs.slides):
            markdown_text += f"\n\n## Slide {i+1}\n\n"
            
            shapes = sorted(slide.shapes, key=lambda x: x.top)

            for shape in shapes:
                if shape.has_text_frame:
                    if shape.is_placeholder and shape.placeholder_format.type in [1, 13, 14, 15, 16]:
                         if shape.text:
                            markdown_text += f"# {shape.text}\n\n"
                    else:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = "".join(run.text for run in paragraph.runs)
                            text = para_text.strip()
                            if not text:
                                continue
                            
                            prefix = ""
                            if paragraph.level > 0:
                                prefix = "  " * paragraph.level + "* "
                            elif text.startswith(('•', '–', '-')):
                                prefix = "* "
                                text = text[1:].lstrip()

                            markdown_text += f"{prefix}{text}\n"
                        markdown_text += "\n"

                elif hasattr(shape, 'image'):
                    markdown_text += f"![Image]\n\n"

        return markdown_text.strip()
    except Exception as e:
        return f"Error converting {params.file_path} to Markdown: {e}"

if __name__ == "__main__":
    mcp.run()
