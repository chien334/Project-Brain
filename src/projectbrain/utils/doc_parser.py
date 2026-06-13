import io
import os
import pypdf
import mammoth
import markdownify

async def call_ocr_if_available(image_bytes: bytes, mime_type: str) -> str:
    """Helper to call vision MCP tool to transcribe image content if credentials are set."""
    try:
        from extensions_mcp.image_to_markdown.vision_client import extract_text_from_image, DEFAULT_PROMPT
        from extensions_mcp.image_to_markdown.config import CONFIG
        
        # If API key is not present, skip OCR call and fallback to standard placeholder
        if not getattr(CONFIG, "api_key", None):
            return "![Image]"
            
        ocr_text = await extract_text_from_image(
            image_bytes=image_bytes,
            mime_type=mime_type,
            prompt=DEFAULT_PROMPT,
            model=CONFIG.model
        )
        if ocr_text:
            return f"\n\n<!-- Start Embedded Image OCR -->\n{ocr_text.strip()}\n<!-- End Embedded Image OCR -->\n\n"
    except Exception:
        pass
    return "![Image]"

async def pdf_to_markdown(file_bytes: bytes) -> str:
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            text_parts = []
            for page_num, page in enumerate(pdf.pages, 1):
                text_parts.append(f"## Page {page_num}\n")
                
                # Extract text
                text = page.extract_text(x_tolerance=2)
                if text:
                    text_parts.append(text + "\n")
                    
                # Extract and format tables
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        # Convert table to markdown
                        if table and len(table) > 0:
                            col_count = max((len(row) for row in table), default=0)
                            if col_count > 0:
                                def normalize_row(row):
                                    cells = [str(row[i] if i < len(row) and row[i] is not None else "").replace("|", "\\|").replace("\n", " ") for i in range(col_count)]
                                    return "| " + " | ".join(cells) + " |"
                                header = normalize_row(table[0])
                                separator = "| " + " | ".join(["---"] * col_count) + " |"
                                body_rows = [normalize_row(row) for row in table[1:]]
                                body = "\n".join(body_rows) if body_rows else ""
                                text_parts.append(f"\n{header}\n{separator}\n{body}\n")
                
                # Extract images and perform OCR if possible
                if hasattr(page, "images") and page.images:
                    for img in page.images:
                        if img.get("width", 0) > 0 and img.get("height", 0) > 0:
                            try:
                                x0 = img.get("x0")
                                top = img.get("top")
                                x1 = img.get("x1")
                                bottom = img.get("bottom")
                                if None not in (x0, top, x1, bottom) and x1 > x0 and bottom > top:
                                    cropped = page.crop((x0, top, x1, bottom))
                                    # Convert cropped to PIL image
                                    img_obj = cropped.to_image(resolution=150)
                                    pil_img = img_obj.original
                                    img_buf = io.BytesIO()
                                    pil_img.save(img_buf, format="PNG")
                                    img_bytes = img_buf.getvalue()
                                    
                                    ocr_result = await call_ocr_if_available(img_bytes, "image/png")
                                    text_parts.append(ocr_result + "\n")
                            except Exception:
                                text_parts.append("![Image]\n")
                                
            return "\n".join(text_parts)
    except ImportError:
        # Fallback to pypdf if pdfplumber is not installed
        try:
            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            text_parts = []
            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                page_text = []
                if text:
                    page_text.append(text)
                
                # Extract images in pypdf if available
                if hasattr(page, "images") and page.images:
                    for img_file in page.images:
                        try:
                            img_bytes = getattr(img_file, "data", None)
                            if img_bytes:
                                name = getattr(img_file, "name", "image.png")
                                mime_type = "image/jpeg" if name.lower().endswith((".jpg", ".jpeg")) else "image/png"
                                ocr_result = await call_ocr_if_available(img_bytes, mime_type)
                                page_text.append(ocr_result)
                        except Exception:
                            page_text.append("![Image]\n")
                
                if page_text:
                    text_parts.append(f"## Page {page_num + 1}\n\n" + "\n".join(page_text))
            return "\n\n".join(text_parts)
        except Exception as e:
            raise ValueError(f"Failed to parse PDF: {str(e)}")
    except Exception as e:
        raise ValueError(f"Failed to parse PDF: {str(e)}")

async def docx_to_markdown(file_bytes: bytes) -> str:
    try:
        file_obj = io.BytesIO(file_bytes)
        result = mammoth.convert_to_html(file_obj)
        html = result.value
        markdown = markdownify.markdownify(html)
        return markdown
    except Exception as e:
        raise ValueError(f"Failed to parse DOCX: {str(e)}")

async def excel_to_markdown(file_bytes: bytes) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        sheets = []
        for name in wb.sheetnames:
            ws = wb[name]
            sheets.append(f"## Sheet: {name}\n")
            all_rows = list(ws.iter_rows(values_only=True))
            max_cols = 0
            # Pre-scan rows to find max columns
            for row in all_rows:
                if row and any(row[i] is not None for i in range(len(row))):
                    max_cols = max(max_cols, len(row))
                    
            if max_cols == 0:
                sheets.append("_Empty sheet_\n")
                continue
                
            first_row = True
            for row in all_rows:
                if not row or not any(row[i] is not None for i in range(len(row))):
                    continue  # skip entirely empty rows
                # Pad row values
                padded_row = [str(row[i]) if i < len(row) and row[i] is not None else "" for i in range(max_cols)]
                cleaned_row = [val.replace("|", "\\|").replace("\n", " ") for val in padded_row]
                row_str = "| " + " | ".join(cleaned_row) + " |"
                
                if first_row:
                    sheets.append(row_str)
                    # Separator line
                    separator = "| " + " | ".join(["---"] * max_cols) + " |"
                    sheets.append(separator)
                    first_row = False
                else:
                    sheets.append(row_str)
            sheets.append("")
        return "\n".join(sheets)
    except ImportError:
        return "Excel parsing requires 'openpyxl'. Run: pip install openpyxl"
    except Exception as e:
        raise ValueError(f"Failed to parse Excel: {str(e)}")

async def pptx_to_markdown(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        markdown_text = []
        for i, slide in enumerate(prs.slides):
            markdown_text.append(f"\n## Slide {i+1}\n")
            
            # Sort shapes by vertical position
            shapes = sorted(slide.shapes, key=lambda x: x.top)
            for shape in shapes:
                if shape.has_text_frame:
                    if shape.is_placeholder and shape.placeholder_format.type in [1, 13, 14, 15, 16]:
                        if shape.text:
                            markdown_text.append(f"# {shape.text.strip()}\n")
                    else:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = "".join(run.text for run in paragraph.runs).strip()
                            if not para_text:
                                continue
                            prefix = ""
                            if paragraph.level > 0:
                                prefix = "  " * paragraph.level + "* "
                            elif para_text.startswith(('•', '–', '-')):
                                prefix = "* "
                                para_text = para_text[1:].lstrip()
                            markdown_text.append(f"{prefix}{para_text}")
                        markdown_text.append("")
                elif hasattr(shape, 'image'):
                    try:
                        img_bytes = shape.image.blob
                        mime_type = shape.image.content_type
                        ocr_result = await call_ocr_if_available(img_bytes, mime_type)
                        markdown_text.append(ocr_result)
                    except Exception:
                        markdown_text.append("![Image]\n")
        return "\n".join(markdown_text).strip()
    except ImportError:
        return "PowerPoint parsing requires 'python-pptx'. Run: pip install python-pptx"
    except Exception as e:
        raise ValueError(f"Failed to parse PPTX: {str(e)}")

async def parse_document(filename: str, file_bytes: bytes) -> str:
    _, ext = os.path.splitext(filename.lower())
    if ext == ".pdf":
        return await pdf_to_markdown(file_bytes)
    elif ext == ".docx":
        return await docx_to_markdown(file_bytes)
    elif ext in [".xlsx", ".xls"]:
        return await excel_to_markdown(file_bytes)
    elif ext == ".pptx":
        return await pptx_to_markdown(file_bytes)
    elif ext in [".txt", ".md", ".json", ".yml", ".yaml", ".ini", ".conf"]:
        return file_bytes.decode("utf-8", errors="ignore")
    else:
        return file_bytes.decode("utf-8", errors="ignore")
